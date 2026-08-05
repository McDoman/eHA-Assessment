"""Shared visual theme for branded DOCX and PPTX deliverables.

A two-colour system. Everything is carried by one blue, one pale tint of it,
black text and white ground - no third hue, no gradients, no shadows.
Hierarchy comes from weight, size and the pale band, never from a new colour.

    slide size      10 x 5.625 in  (9144000 x 5143500 EMU, 16:9)
    primary         #0090FC   titles, rules, stat numbers, section numerals
    pale tint       #CCE9FE   tile fills, table header bands, callouts
    black           #000000   body text
    white           #FFFFFF   ground, reversed text on the primary

Fonts resolve at import time against what is actually installed, so the PPTX
and the PDF exported from it render identically. Override the palette or the
font preference before building:

    import doctheme as T
    T.set_palette("C0392B", "FDEDEC")
    T.set_fonts(["Inter", "Segoe UI", "Arial"])

Requires: python-docx, python-pptx.
"""

from __future__ import annotations

import glob
import os
import re as _re

from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DInches
from docx.shared import Pt as DPt
from docx.shared import RGBColor as DRGB
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# Palette
# --------------------------------------------------------------------------

BLUE = RGBColor(0x00, 0x90, 0xFC)
PALE = RGBColor(0xCC, 0xE9, 0xFE)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEX_BLUE = "0090FC"
HEX_PALE = "CCE9FE"
HEX_BLACK = "000000"
HEX_WHITE = "FFFFFF"

D_BLUE = DRGB(0x00, 0x90, 0xFC)
D_PALE = DRGB(0xCC, 0xE9, 0xFE)
D_BLACK = DRGB(0x00, 0x00, 0x00)
D_WHITE = DRGB(0xFF, 0xFF, 0xFF)


def set_palette(primary_hex: str, pale_hex: str) -> None:
    """Re-brand to a different primary and its pale tint. Call before building.

    Pick the tint at roughly 12-15 per cent of the primary so that black body
    text still clears 4.5:1 against it.
    """
    global BLUE, PALE, HEX_BLUE, HEX_PALE, D_BLUE, D_PALE
    primary_hex = primary_hex.lstrip("#").upper()
    pale_hex = pale_hex.lstrip("#").upper()
    rgb = lambda h: (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    BLUE, PALE = RGBColor(*rgb(primary_hex)), RGBColor(*rgb(pale_hex))
    D_BLUE, D_PALE = DRGB(*rgb(primary_hex)), DRGB(*rgb(pale_hex))
    HEX_BLUE, HEX_PALE = primary_hex, pale_hex


# --------------------------------------------------------------------------
# Fonts
# --------------------------------------------------------------------------

_PREFERRED = ["Open Sans", "Arial"]


def _installed_font_files() -> set[str]:
    names: set[str] = set()
    dirs = [os.path.join(os.environ.get("WINDIR", r"C:\Windows"), "Fonts")]
    local = os.environ.get("LOCALAPPDATA")
    if local:
        dirs.append(os.path.join(local, "Microsoft", "Windows", "Fonts"))
    for d in dirs:
        for ext in ("*.ttf", "*.otf", "*.ttc"):
            for f in glob.glob(os.path.join(d, ext)):
                names.add(os.path.basename(f).lower().replace("-", "").replace("_", ""))
    return names


def resolve_font() -> str:
    """Return the first preferred font family actually installed."""
    files = _installed_font_files()
    for fam in _PREFERRED:
        probe = fam.lower().replace(" ", "")
        if any(probe in f for f in files):
            return fam
    return "Arial"


FONT = resolve_font()
FONT_LIGHT = FONT          # if the family has no Light cut, size/colour carry it
FONT_SEMI = FONT           # emphasis is carried by bold
MONO = "Consolas"          # file names, code, formulas, menu paths


def set_fonts(preferred: list[str]) -> None:
    """Change the font preference order and re-resolve. Call before building.

    Only families actually installed are considered, because an uninstalled
    font silently substitutes at PDF export and the output stops matching.
    """
    global _PREFERRED, FONT, FONT_LIGHT, FONT_SEMI
    _PREFERRED = list(preferred)
    FONT = resolve_font()
    FONT_LIGHT = FONT_SEMI = FONT

# --------------------------------------------------------------------------
# Slide geometry (inches)
# --------------------------------------------------------------------------

SLIDE_W = 10.0
SLIDE_H = 5.625
MARGIN = 0.55
CONTENT_W = SLIDE_W - 2 * MARGIN          # 8.90
GUTTER = 0.34
COL_W = (CONTENT_W - GUTTER) / 2          # 4.28
TITLE_Y = 0.40
TITLE_H = 0.62
BODY_Y = 1.28

# Type scale tuned for a 10-inch canvas. Sized so the densest slide in the
# deck - a five-session day list in two columns - fills roughly 85 per cent of
# the available height, matching the density of the source deck.
SZ_TITLE = 25
SZ_SUBTITLE = 12
SZ_KICKER = 10
SZ_BODY = 11.5
SZ_SMALL = 10
SZ_TINY = 8.6
SZ_BIGNUM = 42
SZ_NUMERAL = 60


# --------------------------------------------------------------------------
# Low-level PPTX helpers
# --------------------------------------------------------------------------

def _tf(shape, wrap=True, margins=(0, 0, 0, 0), anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = wrap
    l, t, r, b = margins
    tf.margin_left, tf.margin_top = Inches(l), Inches(t)
    tf.margin_right, tf.margin_bottom = Inches(r), Inches(b)
    tf.vertical_anchor = anchor
    return tf


def textbox(slide, x, y, w, h, **kw):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return _tf(box, **kw), box


def para(tf, text, size=SZ_BODY, bold=False, color=BLACK, space_after=5,
         space_before=0, line=1.22, align=PP_ALIGN.LEFT, font=None, first=False,
         italic=False):
    """Append a paragraph (or fill the first, empty one) and return it."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font or FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return p


def hang(p, left_in, hanging_in):
    """Give a paragraph a hanging indent so wrapped lines clear the bullet."""
    from pptx.util import Emu as _Emu
    pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
    pPr.set("marL", str(int(_Emu(Inches(left_in)))))
    pPr.set("indent", str(int(-_Emu(Inches(hanging_in)))))
    return p


def rich(tf, segments, size=SZ_BODY, space_after=5, line=1.22, first=False,
         align=PP_ALIGN.LEFT):
    """Append a paragraph built from (text, bold, colour) segments."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line
    for seg in segments:
        text, bold, color = (list(seg) + [False, BLACK])[:3]
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return p


def rect(slide, x, y, w, h, fill=PALE, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def rule(slide, x, y, w, thickness=0.022, color=BLUE):
    return rect(slide, x, y, w, thickness, fill=color)


# --------------------------------------------------------------------------
# Deck: slide factories reproducing the source deck's layouts
# --------------------------------------------------------------------------

class Deck:
    def __init__(self, presentation):
        self.prs = presentation
        self.prs.slide_width = Emu(9144000)
        self.prs.slide_height = Emu(5143500)
        self._blank = self.prs.slide_layouts[6]

    def blank(self):
        return self.prs.slides.add_slide(self._blank)

    # -- title -------------------------------------------------------------
    def title_slide(self, title, subtitle, meta_lines):
        s = self.blank()
        rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
        rect(s, 0, 0, 0.14, SLIDE_H, fill=BLUE)
        tf, _ = textbox(s, MARGIN + 0.25, 1.28, CONTENT_W - 0.4, 2.4)
        para(tf, title, size=30, bold=True, color=BLUE, space_after=10,
             line=1.1, first=True)
        para(tf, subtitle, size=12, color=BLACK, space_after=0, line=1.35)
        tf2, _ = textbox(s, MARGIN + 0.25, 4.05, CONTENT_W - 0.4, 1.1)
        rule(s, MARGIN + 0.25, 3.92, 1.5)
        for i, line in enumerate(meta_lines):
            para(tf2, line, size=SZ_KICKER, color=BLACK, space_after=3,
                 first=(i == 0))
        return s

    # -- section divider ---------------------------------------------------
    def section(self, numeral, title, subtitle=""):
        s = self.blank()
        rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=BLUE)
        tf, _ = textbox(s, MARGIN, 1.15, 3.0, 1.5)
        para(tf, numeral, size=SZ_NUMERAL, bold=True, color=WHITE,
             space_after=0, line=1.0, first=True)
        tf2, _ = textbox(s, MARGIN, 2.75, CONTENT_W, 1.6)
        para(tf2, title, size=25, bold=True, color=WHITE, space_after=8,
             line=1.12, first=True)
        if subtitle:
            para(tf2, subtitle, size=11.5, color=PALE, space_after=0, line=1.3)
        return s

    # -- plain title + optional kicker -------------------------------------
    def _head(self, s, title, kicker=""):
        # A title long enough to wrap would otherwise run through the rule.
        # Arial bold is roughly 0.0077 in per point per character.
        size = SZ_TITLE
        per_line = CONTENT_W / (0.0077 * size)
        if len(title) > per_line:
            size = 19.5
            per_line = CONTENT_W / (0.0077 * size)
        lines = max(1, -(-len(title) // int(per_line)))
        h = 0.30 * lines + 0.30

        tf, _ = textbox(s, MARGIN, TITLE_Y, CONTENT_W, h)
        para(tf, title, size=size, bold=True, color=BLUE, space_after=0,
             line=1.08, first=True)
        y = TITLE_Y + h + 0.06
        rule(s, MARGIN, y, CONTENT_W)
        if kicker:
            ktf, _ = textbox(s, MARGIN, y + 0.10, CONTENT_W, 0.44)
            para(ktf, kicker, size=SZ_KICKER, color=BLACK, space_after=0,
                 line=1.25, first=True)
            return y + 0.62
        return y + 0.20

    # -- two-column body ---------------------------------------------------
    def two_col(self, title, left, right, kicker="", size=None):
        """left/right hold (label, body) or (body,) tuples.

        `body` may be a string, or a list of strings which is rendered as a
        real bullet list with a hanging indent.
        """
        s = self.blank()
        top = self._head(s, title, kicker)
        h = SLIDE_H - top - 0.30
        bsz = size or SZ_SMALL
        for col, x in ((left, MARGIN), (right, MARGIN + COL_W + GUTTER)):
            tf, _ = textbox(s, x, top, COL_W, h)
            firstp = True

            def emit_body(body):
                nonlocal firstp
                if isinstance(body, (list, tuple)):
                    for b in body:
                        p = rich(tf, [("•   ", False, BLUE), (b, False, BLACK)],
                                 size=bsz, space_after=5, line=1.25,
                                 first=firstp)
                        hang(p, 0.24, 0.24)
                        firstp = False
                elif body:
                    para(tf, body, size=bsz, color=BLACK, space_after=4,
                         line=1.25, first=firstp)
                    firstp = False

            for item in col:
                if len(item) == 2:
                    label, body = item
                    if label:
                        para(tf, label, size=bsz + 1, bold=True, color=BLUE,
                             space_after=3, space_before=0 if firstp else 8,
                             first=firstp)
                        firstp = False
                    emit_body(body)
                else:
                    emit_body(item[0])
        return s

    # -- big-number stat tiles (slide 4 idiom) ------------------------------
    def stat_tiles(self, title, tiles, kicker=""):
        s = self.blank()
        top = self._head(s, title, kicker)
        n = len(tiles)
        gap = 0.26
        w = (CONTENT_W - gap * (n - 1)) / n
        h = min(2.55, SLIDE_H - top - 0.35)
        for i, (num, body) in enumerate(tiles):
            x = MARGIN + i * (w + gap)
            rect(s, x, top, w, h, fill=PALE)
            rect(s, x, top, w, 0.055, fill=BLUE)
            tf, _ = textbox(s, x + 0.20, top + 0.26, w - 0.40, h - 0.42)
            para(tf, num, size=SZ_BIGNUM, bold=True, color=BLUE,
                 space_after=7, line=1.0, first=True)
            para(tf, body, size=SZ_SMALL, color=BLACK, space_after=0, line=1.3)
        return s

    # -- bullet list -------------------------------------------------------
    def bullets(self, title, items, kicker="", size=SZ_BODY, two_col=False):
        s = self.blank()
        top = self._head(s, title, kicker)
        h = SLIDE_H - top - 0.30
        if two_col:
            half = (len(items) + 1) // 2
            cols = [items[:half], items[half:]]
            xs = [MARGIN, MARGIN + COL_W + GUTTER]
            ws = [COL_W, COL_W]
        else:
            cols, xs, ws = [items], [MARGIN], [CONTENT_W]
        for col, x, w in zip(cols, xs, ws):
            tf, _ = textbox(s, x, top, w, h)
            for i, it in enumerate(col):
                p = rich(tf, [("\u2022   ", False, BLUE), (it, False, BLACK)],
                         size=size, space_after=7, line=1.25, first=(i == 0))
                hang(p, 0.24, 0.24)
        return s

    # -- session list (Day 1/2/3 idiom, slides 15-17) -----------------------
    def sessions(self, title, kicker, sessions):
        """sessions: list of (code_and_time, content, checkpoint)."""
        s = self.blank()
        top = self._head(s, title, kicker)
        h = SLIDE_H - top - 0.28
        half = (len(sessions) + 1) // 2
        for col, x in ((sessions[:half], MARGIN),
                       (sessions[half:], MARGIN + COL_W + GUTTER)):
            tf, _ = textbox(s, x, top, COL_W, h)
            for i, (code, content, check) in enumerate(col):
                para(tf, code, size=SZ_SMALL, bold=True, color=BLUE,
                     space_after=1, space_before=0 if i == 0 else 7,
                     first=(i == 0))
                para(tf, content, size=SZ_TINY, color=BLACK, space_after=1,
                     line=1.22)
                rich(tf, [("Checkpoint.  ", True, BLACK), (check, False, BLACK)],
                     size=SZ_TINY, space_after=0, line=1.22)
        return s

    # -- table -------------------------------------------------------------
    def table(self, title, header, rows, kicker="", footnote="",
              col_widths=None, size=SZ_TINY, bold_first=None, tight=False):
        """PowerPoint grows table rows to fit their text regardless of the
        height we request, so `tight` trims cell padding for long tables."""
        s = self.blank()
        top = self._head(s, title, kicker)
        avail_h = SLIDE_H - top - (0.78 if footnote else 0.28)
        nrows, ncols = len(rows) + 1, len(header)
        gt = s.shapes.add_table(nrows, ncols, Inches(MARGIN), Inches(top),
                                Inches(CONTENT_W), Inches(avail_h))
        tbl = gt.table
        tbl.first_row = True
        if col_widths:
            total = sum(col_widths)
            for i, cw in enumerate(col_widths):
                tbl.columns[i].width = Emu(int(Inches(CONTENT_W) * cw / total))
        pad = 0.012 if tight else 0.03
        hdr_h = Inches(0.26 if tight else 0.30)
        tbl.rows[0].height = hdr_h
        body_h = Emu(max(int((Inches(avail_h) - hdr_h) / max(len(rows), 1)),
                         int(Inches(0.16))))
        for i in range(1, nrows):
            tbl.rows[i].height = body_h
        if bold_first is None:
            bold_first = ncols > 2
        for c, text in enumerate(header):
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
            self._cell(cell, text, size=size, bold=True, color=WHITE, pad=pad)
        for r, row in enumerate(rows, start=1):
            for c, text in enumerate(row):
                cell = tbl.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PALE if r % 2 else WHITE
                self._cell(cell, str(text), size=size,
                           bold=(c == 0 and bold_first), pad=pad)
        if footnote:
            tf, _ = textbox(s, MARGIN, SLIDE_H - 0.72, CONTENT_W, 0.58)
            para(tf, footnote, size=SZ_TINY, color=BLACK, space_after=0,
                 line=1.2, first=True)
        return s

    @staticmethod
    def _cell(cell, text, size=SZ_TINY, bold=False, color=BLACK, pad=0.03):
        cell.margin_left = Inches(0.07)
        cell.margin_right = Inches(0.07)
        cell.margin_top = Inches(pad)
        cell.margin_bottom = Inches(pad)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.12
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    # -- full-bleed statement (slide 30 idiom) ------------------------------
    def statement(self, headline, body_lines, footer=""):
        s = self.blank()
        rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=PALE)
        rect(s, 0, 0, SLIDE_W, 0.10, fill=BLUE)
        tf, _ = textbox(s, MARGIN + 0.3, 0.85, CONTENT_W - 0.6, 1.3)
        para(tf, headline, size=21, bold=True, color=BLUE, space_after=0,
             line=1.14, first=True)
        tf2, _ = textbox(s, MARGIN + 0.3, 2.35, CONTENT_W - 0.6, 2.2)
        for i, line in enumerate(body_lines):
            para(tf2, line, size=SZ_BODY, color=BLACK, space_after=7,
                 line=1.3, first=(i == 0))
        if footer:
            tf3, _ = textbox(s, MARGIN + 0.3, SLIDE_H - 0.62, CONTENT_W - 0.6, 0.4)
            para(tf3, footer, size=SZ_KICKER, color=BLACK, space_after=0,
                 first=True)
        return s

    # -- click-path slide for the manual deck ------------------------------
    def steps(self, title, kicker, steps, checkpoint="", watch="", lead=""):
        s = self.blank()
        top = self._head(s, title, kicker)
        right_w = 3.05
        left_w = CONTENT_W - right_w - GUTTER
        h = SLIDE_H - top - 0.28
        tf, _ = textbox(s, MARGIN, top, left_w, h)
        firstp = True
        if lead:
            para(tf, lead, size=SZ_TINY, color=BLACK, space_after=9,
                 line=1.28, first=True)
            firstp = False
        for i, st in enumerate(steps):
            p = rich(tf, [(f"{i + 1}.  ", True, BLUE), (st, False, BLACK)],
                     size=SZ_SMALL, space_after=6, line=1.25,
                     first=(firstp and i == 0))
            hang(p, 0.26, 0.26)
        x = MARGIN + left_w + GUTTER
        if checkpoint:
            ch = 1.30
            rect(s, x, top, right_w, ch, fill=PALE)
            rect(s, x, top, 0.045, ch, fill=BLUE)
            ctf, _ = textbox(s, x + 0.18, top + 0.14, right_w - 0.32, ch - 0.24)
            para(ctf, "Checkpoint", size=SZ_KICKER, bold=True, color=BLUE,
                 space_after=4, first=True)
            para(ctf, checkpoint, size=SZ_TINY, color=BLACK, space_after=0,
                 line=1.25)
            if watch:
                wtf, _ = textbox(s, x, top + ch + 0.18, right_w, h - ch - 0.20)
                para(wtf, "Watch out for", size=SZ_KICKER, bold=True,
                     color=BLUE, space_after=4, first=True)
                para(wtf, watch, size=SZ_TINY, color=BLACK, space_after=0,
                     line=1.25)
        return s


# --------------------------------------------------------------------------
# DOCX helpers
# --------------------------------------------------------------------------

def docx_setup(doc, margin=0.85):
    """Apply the theme to a python-docx Document and set A4 portrait."""
    from docx.enum.section import WD_ORIENT

    for section in doc.sections:
        section.page_width = DInches(8.27)
        section.page_height = DInches(11.69)
        section.orientation = WD_ORIENT.PORTRAIT
        section.left_margin = DInches(margin)
        section.right_margin = DInches(margin)
        section.top_margin = DInches(margin)
        section.bottom_margin = DInches(margin)

    normal = doc.styles["Normal"]
    normal.font.name = FONT
    normal.font.size = DPt(10)
    normal.font.color.rgb = D_BLACK
    normal.paragraph_format.space_after = DPt(6)
    normal.paragraph_format.line_spacing = 1.22
    rpr = normal.element.get_or_add_rPr()
    rfonts = rpr.get_or_add_rFonts()
    for attr in ("w:ascii", "w:hAnsi", "w:cs", "w:eastAsia"):
        rfonts.set(qn(attr), FONT)

    specs = {
        "Title": (24, True, D_BLUE, 0, 6),
        "Heading 1": (16, True, D_BLUE, 16, 6),
        "Heading 2": (12.5, True, D_BLACK, 12, 4),
        "Heading 3": (10.5, True, D_BLUE, 9, 3),
    }
    for name, (size, bold, color, before, after) in specs.items():
        try:
            st = doc.styles[name]
        except KeyError:
            continue
        st.font.name = FONT
        st.font.size = DPt(size)
        st.font.bold = bold
        st.font.color.rgb = color
        st.paragraph_format.space_before = DPt(before)
        st.paragraph_format.space_after = DPt(after)
        st.paragraph_format.keep_with_next = True
        r = st.element.get_or_add_rPr().get_or_add_rFonts()
        for attr in ("w:ascii", "w:hAnsi", "w:cs"):
            r.set(qn(attr), FONT)
    return doc


def shade(cell, hex_fill):
    el = OxmlElement("w:shd")
    el.set(qn("w:val"), "clear")
    el.set(qn("w:color"), "auto")
    el.set(qn("w:fill"), hex_fill)
    cell._tc.get_or_add_tcPr().append(el)


def cell_text(cell, text, size=9, bold=False, color=D_BLACK, align=None):
    cell.text = ""
    p = cell.paragraphs[0]
    p.paragraph_format.space_after = DPt(2)
    p.paragraph_format.space_before = DPt(2)
    p.paragraph_format.line_spacing = 1.15
    if align is not None:
        p.alignment = align
    run = p.add_run(str(text))
    run.font.name = FONT
    run.font.size = DPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return cell


def add_table(doc, header, rows, widths=None, size=9, header_fill=HEX_BLUE,
              zebra=True, bold_first=None):
    """Banded table. Pass `header=None` for a label/value block.

    A label/value table - the usual "At a glance" opener - has no column
    headings, and emitting an empty primary bar above it just looks like a
    mistake. `header=None` (or a header whose cells are all blank) drops the
    header row and bolds the first column instead.
    """
    if header is not None and not any(str(h).strip() for h in header):
        header = None
    ncols = len(header) if header else len(rows[0])
    if bold_first is None:
        bold_first = header is None or ncols > 2

    offset = 1 if header else 0
    t = doc.add_table(rows=len(rows) + offset, cols=ncols)
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    t.style = "Table Grid"
    if header:
        for c, h in enumerate(header):
            shade(t.rows[0].cells[c], header_fill)
            cell_text(t.rows[0].cells[c], h, size=size, bold=True,
                      color=D_WHITE)
    for r, row in enumerate(rows, start=offset):
        for c, val in enumerate(row):
            if zebra and (r - offset) % 2 == 0:
                shade(t.rows[r].cells[c], HEX_PALE)
            cell_text(t.rows[r].cells[c], val, size=size,
                      bold=(c == 0 and bold_first))
    if widths:
        total = sum(widths)
        usable = doc.sections[0].page_width - doc.sections[0].left_margin \
            - doc.sections[0].right_margin
        t.autofit = False
        for r in t.rows:
            for c, w in enumerate(widths):
                r.cells[c].width = int(usable * w / total)
    doc.add_paragraph()
    return t


def callout(doc, kind, body, fill=HEX_PALE):
    """A single-cell shaded box - the draft manual's 'Check your work' idiom."""
    t = doc.add_table(rows=1, cols=1)
    t.style = "Table Grid"
    cell = t.rows[0].cells[0]
    shade(cell, fill)
    cell.text = ""
    p = cell.paragraphs[0]
    p.paragraph_format.space_after = DPt(2)
    r = p.add_run(kind)
    r.font.name, r.font.size, r.font.bold = FONT, DPt(9.5), True
    r.font.color.rgb = D_BLUE
    p2 = cell.add_paragraph()
    p2.paragraph_format.space_after = DPt(2)
    p2.paragraph_format.line_spacing = 1.2
    r2 = p2.add_run(body)
    r2.font.name, r2.font.size = FONT, DPt(9)
    r2.font.color.rgb = D_BLACK
    doc.add_paragraph()
    return t


def bullet(doc, text, level=0, bold_lead=None):
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.left_indent = DInches(0.25 + 0.22 * level)
    p.paragraph_format.space_after = DPt(3)
    p.paragraph_format.line_spacing = 1.2
    if bold_lead:
        r = p.add_run(bold_lead)
        r.font.name, r.font.size, r.font.bold = FONT, DPt(10), True
        r.font.color.rgb = D_BLACK
    r = p.add_run(text)
    r.font.name, r.font.size = FONT, DPt(10)
    r.font.color.rgb = D_BLACK
    return p


def numbered(doc, text, bold_lead=None):
    p = doc.add_paragraph(style="List Number")
    p.paragraph_format.left_indent = DInches(0.28)
    p.paragraph_format.space_after = DPt(3)
    p.paragraph_format.line_spacing = 1.2
    if bold_lead:
        r = p.add_run(bold_lead)
        r.font.name, r.font.size, r.font.bold = FONT, DPt(10), True
    r = p.add_run(text)
    r.font.name, r.font.size = FONT, DPt(10)
    return p


def kicker_para(doc, text):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = DPt(8)
    r = p.add_run(text)
    r.font.name, r.font.size, r.font.italic = FONT, DPt(9), False
    r.font.color.rgb = D_BLUE
    r.font.bold = True
    return p


def hrule(doc):
    p = doc.add_paragraph()
    pr = p._p.get_or_add_pPr()
    bd = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "12")
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), HEX_BLUE)
    bd.append(bottom)
    pr.append(bd)
    return p


def cover(doc, title, subtitle="", meta_lines=(), kicker=""):
    """Standard first block: big primary title, subtitle, rule, meta lines."""
    p = doc.add_paragraph()
    p.paragraph_format.space_after = DPt(2)
    r = p.add_run(title)
    r.font.name, r.font.size, r.font.bold = FONT, DPt(24), True
    r.font.color.rgb = D_BLUE
    if subtitle:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = DPt(2)
        r = p.add_run(subtitle)
        r.font.name, r.font.size = FONT, DPt(12.5)
    if kicker:
        p = doc.add_paragraph()
        r = p.add_run(kicker)
        r.font.name, r.font.size, r.font.bold = FONT, DPt(10), True
        r.font.color.rgb = D_BLUE
    hrule(doc)
    for line in meta_lines:
        p = doc.add_paragraph(line)
        p.paragraph_format.space_after = DPt(2)
        for r in p.runs:
            r.font.name, r.font.size = FONT, DPt(9)
    return doc


# -- inline markup: **bold** and `code` -------------------------------------

_TOKEN = _re.compile(r"(\*\*.+?\*\*|`.+?`)")


def rich_run(p, text, size=10, base_bold=False):
    """Fill a paragraph from a string using **bold** and `code` markers.

    `code` is set in MONO and the primary colour - use it for file names,
    formulas, field names and menu paths, which is most of what a procedural
    document points at.
    """
    for tok in _TOKEN.split(text):
        if not tok:
            continue
        if tok.startswith("**") and tok.endswith("**"):
            r = p.add_run(tok[2:-2])
            r.font.name, r.font.bold = FONT, True
        elif tok.startswith("`") and tok.endswith("`"):
            r = p.add_run(tok[1:-1])
            r.font.name = MONO
            r.font.color.rgb = D_BLUE
        else:
            r = p.add_run(tok)
            r.font.name, r.font.bold = FONT, base_bold
        r.font.size = DPt(size)
    return p


def rich_bullet(doc, text, level=0, size=10):
    """A bullet whose text supports **bold** and `code`."""
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.left_indent = DInches(0.25 + 0.22 * level)
    p.paragraph_format.space_after = DPt(3)
    p.paragraph_format.line_spacing = 1.2
    return rich_run(p, text, size=size)


def strip_md(text):
    """Drop the markers - for contexts that cannot carry mixed runs."""
    return text.replace("**", "").replace("`", "")
